#!/usr/bin/env python3
"""Build the NOVAFRICA conference presentation from the IFPRI template."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import copy, os

SRC = "presentations/nova_africa.pptx"
OUT = "presentations/nova_africa.pptx"
FIG = "paper/figures/"

# IFPRI palette
GREEN  = RGBColor(0x62, 0xBA, 0x45)
TEAL   = RGBColor(0x00, 0xAE, 0x9A)
RED    = RGBColor(0xEF, 0x46, 0x3B)
ORANGE = RGBColor(0xF7, 0x92, 0x1E)
PURPLE = RGBColor(0x88, 0x50, 0xA0)
BLUE   = RGBColor(0x00, 0x7D, 0xB3)
SLATE  = RGBColor(0x44, 0x54, 0x6A)
GREY   = RGBColor(0x6E, 0x6E, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2A, 0x2A, 0x2A)

SW, SH = Inches(13.33), Inches(7.5)

prs = Presentation(SRC)

# ---- layouts ----
L = {l.name: l for l in prs.slide_masters[0].slide_layouts}
LAY_CONTENT = L['Title and Content 1']

# =====================================================================
# helpers
# =====================================================================
def remove_placeholder(slide, idx):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            ph._element.getparent().remove(ph._element)

def set_title(slide, text, size=30, color=SLATE):
    ph = slide.placeholders[0]
    ph.text = text
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = color; r.font.name = 'Arial'
    return ph

def add_text(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tb, tf

def para(tf, text, size=18, color=DARK, bold=False, italic=False, bullet=False,
         align=PP_ALIGN.LEFT, space_after=8, level=0, first=False, color_run=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.level = level
    p.space_after = Pt(space_after)
    segs = text if isinstance(text, list) else [(text, color, bold)]
    for seg in segs:
        txt, col, bd = seg
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.name = 'Arial'
        r.font.bold = bd; r.font.italic = italic
        r.font.color.rgb = col
    if bullet:
        _set_bullet(p)
    else:
        _no_bullet(p)
    return p

from pptx.oxml.ns import qn
def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ('a:buChar','a:buAutoNum'):
        for e in pPr.findall(qn(tag)): pPr.remove(e)
    if pPr.find(qn('a:buNone')) is None:
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))

def _set_bullet(p, char='•', color=None):
    pPr = p._p.get_or_add_pPr()
    for tag in ('a:buNone','a:buChar','a:buAutoNum'):
        for e in pPr.findall(qn(tag)): pPr.remove(e)
    lvl = p.level or 0
    pPr.set('marL', str(int(Inches(0.30) + Inches(0.30)*lvl)))
    pPr.set('indent', str(-int(Inches(0.30))))
    buf = pPr.makeelement(qn('a:buFont'), {'typeface':'Arial'}); pPr.append(buf)
    bu = pPr.makeelement(qn('a:buChar'), {'char':char}); pPr.append(bu)

def img_fit(slide, path, box_l, box_t, box_w, box_h, align='center', valign='middle'):
    im = Image.open(path); iw, ih = im.size; ar = iw/ih
    bw, bh = box_w, box_h
    if bw/bh > ar:
        h = bh; w = bh*ar
    else:
        w = bw; h = bw/ar
    if align=='center': l = box_l + (bw-w)/2
    elif align=='left': l = box_l
    else: l = box_l + (bw-w)
    if valign=='middle': t = box_t + (bh-h)/2
    elif valign=='top': t = box_t
    else: t = box_t + (bh-h)
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def rect(slide, l, t, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s

def content_slide(title, tsize=30):
    s = prs.slides.add_slide(LAY_CONTENT)
    remove_placeholder(s, 1)
    set_title(s, title, size=tsize)
    return s

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def callout(slide, l, t, w, h, big, big_color, label, sub=None):
    r = rect(slide, l, t, w, h, RGBColor(0xF5,0xF7,0xF9), line=RGBColor(0xD5,0xDB,0xE0))
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_top=Pt(6); tf.margin_bottom=Pt(6)
    para(tf, big, size=30, color=big_color, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
    para(tf, label, size=13, color=SLATE, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    if sub:
        para(tf, sub, size=11, color=GREY, align=PP_ALIGN.CENTER, space_after=0)
    return r

# =====================================================================
# SLIDE 1 — Title (reuse Alt Title Page = slide index 1)
# =====================================================================
title = prs.slides[1]
def setruns(name, lines):
    sh = [s for s in title.shapes if s.name==name][0]
    runs = sh.text_frame.paragraphs[0].runs
    # ensure enough paragraphs for multi-line title boxes handled separately
    return sh

# TextBox 1 (two-line title)
tb1 = [s for s in title.shapes if s.name=='TextBox 1'][0]
tb1.text_frame.paragraphs[0].runs[0].text = "When Quality (doesn’t) Pay"
tb1.text_frame.paragraphs[1].runs[0].text = "Two Experiments in Ugandan Dairy"
# TextBox 2 subtitle
tb2 = [s for s in title.shapes if s.name=='TextBox 2'][0]
tb2.text_frame.paragraphs[0].runs[0].text = "Technology, transparency, and the emergence of markets for quality"
# TextBox 3 authors/venue
tb3 = [s for s in title.shapes if s.name=='TextBox 3'][0]
r3 = tb3.text_frame.paragraphs[0].runs
vals = ["Bjorn Van Campenhout",
        "with R. Ariong, S. W. Kariuki & J. Chamberlin",
        "IFPRI  |  KU Leuven",
        "NOVAFRICA Conference  |  19 June 2026"]
for i,r in enumerate(r3):
    r.text = vals[i] if i < len(vals) else ""
notes(title, "Two linked field experiments in Uganda's dairy sector asking a single question: "
      "can making quality observable create a market for quality? "
      "Spoiler: observability is necessary but not sufficient.")

# =====================================================================
# SLIDE 2 — Problem 1: quality is unobservable
# =====================================================================
s = content_slide("The problem (1): quality you cannot see")
_, tf = add_text(s, 0.7, 1.7, 7.4, 5.2)
para(tf, "Milk quality (butterfat, solids-not-fat, added water) is "
        "costly to verify at the point of sale.", size=21, bold=True, color=SLATE, first=True, space_after=14)
para(tf, [("MCCs rely on crude tools: ", DARK, False),
          ("lactometer + alcohol test", BLUE, True),
          (" — detect gross adulteration, not composition.", DARK, False)],
     size=18, bullet=True, space_after=10)
para(tf, [("High- and low-quality milk therefore trade at the ", DARK, False),
          ("same per-litre price", RED, True),
          (".", DARK, False)], size=18, bullet=True, space_after=10)
para(tf, [("Classic ", DARK, False), ("asymmetric information", PURPLE, True),
          (" (Akerlof 1970): no credible signal → no reward for quality → quality erodes.", DARK, False)],
     size=18, bullet=True, space_after=10)
para(tf, [("Both adulteration margins pay: ", DARK, False),
          ("watering", RED, True), (" inflates volume; ", DARK, False),
          ("skimming fat", RED, True), (" extracts the most valuable component.", DARK, False)],
     size=18, bullet=True, space_after=10)
# right pull-quote box
rb = rect(s, 8.5, 2.1, 4.2, 3.2, RGBColor(0xF3,0xF6,0xF2), line=None)
tf2 = rb.text_frame; tf2.word_wrap=True; tf2.vertical_anchor=MSO_ANCHOR.MIDDLE
tf2.margin_left=Pt(14); tf2.margin_right=Pt(14)
para(tf2, "“Without verification, the market cannot tell good milk from bad.”",
     size=20, italic=True, bold=True, color=SLATE, align=PP_ALIGN.CENTER, first=True, space_after=10)
para(tf2, "5–10% of milk samples in the region show water addition", size=14, color=GREY, align=PP_ALIGN.CENTER)
notes(s, "Start with the information problem. Quality is a credence/experience attribute at first sale. "
      "Crude tests catch only gross adulteration. Pay-per-litre with no composition adjustment means both "
      "watering and skimming are profitable. This is the lemons problem.")

# =====================================================================
# SLIDE 3 — Problem 2: aggregation kills the market for quality
# =====================================================================
s = content_slide("The problem (2): aggregation severs the link")
_, tf = add_text(s, 0.7, 1.7, 7.4, 5.2)
para(tf, "Milk is pooled across many farmers early in the chain.", size=21, bold=True, color=SLATE, first=True, space_after=14)
para(tf, [("Buyers observe only a ", DARK, False), ("group-level", TEAL, True),
          (" quality signal; individual contributions vanish into the pool.", DARK, False)],
     size=18, bullet=True, space_after=10)
para(tf, [("Shared reputation = a ", DARK, False), ("common-pool resource", PURPLE, True),
          (" (Winfree & McCluskey 2005; Tirole 1996).", DARK, False)], size=18, bullet=True, space_after=10)
para(tf, [("Moral hazard in teams: each farmer can ", DARK, False),
          ("free-ride", RED, True),
          (" on the pool — no one internalises the return to their own effort.", DARK, False)],
     size=18, bullet=True, space_after=10)
para(tf, [("Layers of traders make ", DARK, False), ("traceability", BLUE, True),
          (" break down even faster.", DARK, False)], size=18, bullet=True, space_after=14)
para(tf, [("⇒  Even if buyers ", SLATE, True), ("want", SLATE, True),
          (" to pay for quality, the structure prevents it. ", SLATE, True),
          ("No market for quality emerges.", RED, True)],
     size=19, space_after=0)
# right: schematic
rb = rect(s, 8.6, 2.0, 4.1, 3.6, RGBColor(0xF3,0xF6,0xF2), line=None)
tf2 = rb.text_frame; tf2.word_wrap=True; tf2.vertical_anchor=MSO_ANCHOR.MIDDLE
tf2.margin_left=Pt(12); tf2.margin_right=Pt(12)
para(tf2, "Farmers  →  Trader  →  MCC (pool)  →  Processor",
     size=15, bold=True, color=SLATE, align=PP_ALIGN.CENTER, first=True, space_after=12)
para(tf2, "individual quality\nobserved with noise", size=13, color=GREY, align=PP_ALIGN.CENTER, space_after=12)
para(tf2, "pooled & paid one\nflat price per litre", size=13, color=GREY, align=PP_ALIGN.CENTER, space_after=0)
notes(s, "Second, conceptually distinct impediment: aggregation. Pooling makes reputation a common-pool resource; "
      "moral-hazard-in-teams logic means individual effort is not rewarded. Intermediation layers worsen traceability. "
      "Two barriers: (i) observability at point of sale, (ii) traceability through aggregation.")

# =====================================================================
# SLIDE 4 — This paper / roadmap
# =====================================================================
s = content_slide("This paper: two linked field experiments")
_, tf = add_text(s, 0.7, 1.6, 12.0, 0.9)
para(tf, [("Can making quality ", SLATE, True), ("observable and traceable", TEAL, True),
          (" create a market for quality?", SLATE, True)], size=22, first=True, space_after=0)
# two big boxes
b1 = rect(s, 0.7, 2.7, 5.7, 3.6, RGBColor(0xEC,0xF5,0xE9), line=GREEN)
tf1 = b1.text_frame; tf1.word_wrap=True; tf1.margin_left=Pt(16); tf1.margin_right=Pt(16); tf1.margin_top=Pt(12)
para(tf1, "Experiment 1", size=20, bold=True, color=GREEN, first=True, space_after=4)
para(tf1, "Measurement & digital monitoring", size=16, bold=True, color=SLATE, space_after=10)
para(tf1, [("Cluster-RCT, 95 catchments. Install milk analyzers + record-keeping app at MCCs.", DARK, False)], size=15, bullet=True, space_after=8)
para(tf1, [("Tests the ", DARK, False), ("information / enforcement", GREEN, True), (" margin.", DARK, False)], size=15, bullet=True, space_after=0)
b2 = rect(s, 6.9, 2.7, 5.7, 3.6, RGBColor(0xFD,0xF0,0xE4), line=ORANGE)
tf2 = b2.text_frame; tf2.word_wrap=True; tf2.margin_left=Pt(16); tf2.margin_right=Pt(16); tf2.margin_top=Pt(12)
para(tf2, "Experiment 2", size=20, bold=True, color=ORANGE, first=True, space_after=4)
para(tf2, "Activating the price channel", size=16, bold=True, color=SLATE, space_after=10)
para(tf2, [("Within-MCC RCT, ~100 traders. Pay traders an explicit quality premium.", DARK, False)], size=15, bullet=True, space_after=8)
para(tf2, [("Tests the ", DARK, False), ("price / pass-through", ORANGE, True), (" margin.", DARK, False)], size=15, bullet=True, space_after=0)
notes(s, "Roadmap. The sequential design separates the informational dimension (Exp 1) from the pricing dimension (Exp 2) "
      "in the same supply chain. Exp 2 is run on infrastructure built by Exp 1.")

# =====================================================================
# SLIDE 5 — Context
# =====================================================================
s = content_slide("Context: Uganda’s dairy supply chain")
_, tf = add_text(s, 0.7, 1.7, 6.2, 5.2)
para(tf, "Southwestern milk shed (Mbarara)", size=19, bold=True, color=SLATE, first=True, space_after=12)
para(tf, [("Most commercialised livestock sub-sector: ", DARK, False), ("17% of agricultural GDP", BLUE, True), (".", DARK,False)], size=17, bullet=True, space_after=8)
para(tf, [("Output ", DARK, False), ("5× since 2008", GREEN, True), ("; exports up to USD 265m (mostly to Kenya).", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("Dense network of ", DARK, False), ("milk collection centres (MCCs)", TEAL, True), (" linking farmers to processors.", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("Quality is central: it sets processing yields for cheese, casein, infant formula.", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("Yet prices are ", DARK, False), ("fixed per litre", RED, True),
          (", set for 2-week windows, paid biweekly — the quality–payment link is severed.", DARK, False)], size=17, bullet=True, space_after=0)
img_fit(s, FIG+"milk_analyzer_installed.jpg", 7.2, 1.7, 5.5, 5.1, align='center', valign='top')
add_text(s, 7.2, 6.85, 5.5, 0.4)[1]
notes(s, "Set the scene: a transforming, commercially organised sector where quality genuinely matters for downstream "
      "products, but pricing institutions are flat per-litre. Photo: milk analyzer installed at an MCC during piloting.")

# =====================================================================
# SLIDE 6 — Experiment 1 design
# =====================================================================
s = content_slide("Experiment 1: design")
_, tf = add_text(s, 0.7, 1.7, 6.2, 5.2)
para(tf, "Cluster-randomised controlled trial", size=19, bold=True, color=GREEN, first=True, space_after=12)
para(tf, [("Unit of randomisation: ", DARK, False), ("catchment area", TEAL, True),
          (" (1–few co-located MCCs + their farmers).", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("95 clusters / 125 MCCs", BLUE, True), ("; ~half treated, half control.", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("20 farmers", BLUE, True), (" sampled per MCC → ~2,500 farmers.", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("Effects estimated at ", DARK, False), ("three levels", PURPLE, True),
          (": MCC, farmer, and 2,518 supervised ", DARK, False), ("milk samples", PURPLE, True),
          (" at endline.", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("Baseline 2022  →  install/train 2023  →  endline Dec 2024.", DARK, False)], size=17, bullet=True, space_after=0)
img_fit(s, FIG+"sampling_two_panels.png", 7.1, 1.9, 5.6, 4.4, align='center', valign='top')
_, cap = add_text(s, 7.1, 6.25, 5.6, 0.6)
para(cap, "Treatment assignment: MCCs (left) and their farmers (right). Green = treated, red = control.",
     size=11, color=GREY, align=PP_ALIGN.CENTER, first=True, space_after=0)
notes(s, "Cluster RCT at the catchment level because the intervention is at the MCC and farmers within a catchment share MCCs. "
      "Three measurement levels including objective supervised milk testing — the cleanest outcome.")

# =====================================================================
# SLIDE 7 — Experiment 1 intervention
# =====================================================================
s = content_slide("Experiment 1: the intervention bundle")
_, tf = add_text(s, 0.7, 1.7, 7.9, 5.2)
para(tf, [("1.  Milk analyzer", GREEN, True), ("  — measures fat, SNF, protein, added water & density in <1 min/sample.",
          DARK, False)], size=18, first=True, space_after=12)
para(tf, [("2.  Digital record-keeping app", TEAL, True),
          ("  — Android tablet replaces paper notebook; stores quality + quantity per farmer, builds quality histories.",
          DARK, False)], size=18, space_after=12)
para(tf, [("3.  Farmer awareness campaign", BLUE, True),
          ("  — posters publicise the new testing capacity and farmers’ right to request a free test.",
          DARK, False)], size=18, space_after=16)
para(tf, [("Together: make quality ", SLATE, True), ("observable, traceable, and actionable", PURPLE, True),
          (".", SLATE, True)], size=18, space_after=10)
para(tf, [("Compliance was imperfect — first stage = ", DARK, False), ("0.37", RED, True),
          (" (24% → 61% with a functioning analyzer). ITT is a lower bound.", DARK, False)], size=15, italic=True, color=GREY, space_after=0)
img_fit(s, FIG+"app.png", 8.9, 1.7, 3.7, 5.1, align='center', valign='top')
notes(s, "Three components: hardware (analyzer), software (app + portals for owners and DDA), and a farmer-facing "
      "information campaign to avoid power imbalance. Note imperfect compliance: analyzers moved, some unused. "
      "First stage 0.37, so ITT attenuated; we also report LATE.")

# =====================================================================
# SLIDE 8 — Experiment 1 results: quality improves
# =====================================================================
s = content_slide("Experiment 1 results: milk quality improves")
_, tf = add_text(s, 0.7, 1.55, 12.0, 0.8)
para(tf, [("On 2,518 supervised milk samples, treated MCCs receive ", SLATE, False),
          ("cleaner, richer milk", GREEN, True), (".", SLATE, False)], size=20, first=True, space_after=0)
y=2.55; w=2.75; h=1.7; gap=0.25
callout(s, 0.7,       y, w, h, "+0.11pp", GREEN,  "Butterfat **",   "3.88% → ~4.0%")
callout(s, 0.7+(w+gap), y, w, h, "−0.49pp", RED, "Added water *",  "vs 1.63% control")
callout(s, 0.7+2*(w+gap), y, w, h, "+0.44", TEAL,  "Density (CLR) *", "less dilution")
callout(s, 0.7+3*(w+gap), y, w, h, "+0.21σ", PURPLE, "Quality index **", "Anderson composite")
_, tf2 = add_text(s, 0.7, 4.6, 12.0, 2.2)
para(tf2, [("Effects load on ", DARK, False), ("handling/adulteration margins", GREEN, True),
           (" (fat, water, density), not on biologically-determined SNF and protein (both null).", DARK, False)],
     size=18, bullet=True, first=True, space_after=10)
para(tf2, [("Pattern is consistent with ", DARK, False), ("reduced skimming and dilution", TEAL, True),
           (", not deeper changes in feeding or production.", DARK, False)], size=18, bullet=True, space_after=10)
para(tf2, [("Robust to dropping ", DARK, False), ("contaminated controls", GREY, True),
           (" (index +0.22σ on clean sample) and to LATE.", DARK, False)], size=18, bullet=True, space_after=0)
notes(s, "Headline of Exp 1: quality improves, but selectively — on the margins a farmer/trader can manipulate quickly "
      "(water, skimming) rather than on composition that needs feeding changes. Significance stars: **=1%, *=5%.")

# =====================================================================
# SLIDE 9 — Experiment 1 results: but no premium (forest plots)
# =====================================================================
s = content_slide("Experiment 1 results: … but no price premium emerges")
img_fit(s, FIG+"fig_forest_mcc_primary.png",   0.4, 1.6, 4.25, 3.4, align='center', valign='top')
img_fit(s, FIG+"fig_forest_farmer_primary.png",4.55,1.6, 4.25, 3.4, align='center', valign='top')
img_fit(s, FIG+"fig_forest_prices_exp1.png",   8.7, 1.6, 4.25, 3.4, align='center', valign='top')
_, tf = add_text(s, 0.7, 5.25, 12.0, 1.9)
para(tf, [("Testing behaviour jumps", GREEN, True),
          (": treated MCCs test incoming milk +53pp (29% → 82%).", DARK, False)], size=17, bullet=True, first=True, space_after=8)
para(tf, [("Prices and premiums: ", DARK, False), ("all null", RED, True),
          (" — farm-gate price, MCC buy/sell price, quality bonuses (MCC & farmer).", DARK, False)], size=17, bullet=True, space_after=8)
para(tf, [("Quality rises, but no one is paid for it. ", SLATE, True),
          ("Enforcement without a price.", PURPLE, True)], size=18, space_after=0)
notes(s, "Three forest plots: MCC primary outcomes, farmer primary outcomes, farmer prices by season. "
      "Testing rises sharply; every price/premium coefficient straddles zero. The quality gain is real but unrewarded — "
      "this is the puzzle the model resolves.")

# =====================================================================
# SLIDE 10 — Model: two channels
# =====================================================================
s = content_slide("A model: two channels for quality infrastructure")
_, tf = add_text(s, 0.7, 1.55, 12.0, 0.8)
para(tf, [("MCC sees a noisy signal ", SLATE, False), ("sᵢ = qᵢ + ηᵢ,  ηᵢ ~ N(0, σ²)", BLUE, True),
          (";  accepts if sᵢ ≥ q_min.  The analyzer drives σ² → 0.", SLATE, False)], size=18, first=True, space_after=0)
b1 = rect(s, 0.7, 2.6, 5.8, 3.9, RGBColor(0xEC,0xF5,0xE9), line=GREEN)
t1 = b1.text_frame; t1.word_wrap=True; t1.margin_left=Pt(16); t1.margin_right=Pt(16); t1.margin_top=Pt(12)
para(t1, "Enforcement channel", size=19, bold=True, color=GREEN, first=True, space_after=8)
para(t1, [("Lower σ² makes the rejection threshold ", DARK, False), ("credibly bind", GREEN, True), (".", DARK, False)], size=15, bullet=True, space_after=8)
para(t1, [("Farmers near q_min raise effort to ", DARK, False), ("avoid rejection", DARK, True),
          (" — the extensive margin.", DARK, False)], size=15, bullet=True, space_after=8)
para(t1, [("Works even if the price premium α = 0.", DARK, False)], size=15, bullet=True, space_after=8)
para(t1, [("Signature: quality ↑ (esp. lower tail), prices flat.", GREEN, True)], size=15, bold=True, space_after=0)
b2 = rect(s, 6.8, 2.6, 5.8, 3.9, RGBColor(0xFD,0xF0,0xE4), line=ORANGE)
t2 = b2.text_frame; t2.word_wrap=True; t2.margin_left=Pt(16); t2.margin_right=Pt(16); t2.margin_top=Pt(12)
para(t2, "Price channel", size=19, bold=True, color=ORANGE, first=True, space_after=8)
para(t2, [("Requires a per-unit quality premium ", DARK, False), ("α > 0", ORANGE, True),
          (" paid up the chain.", DARK, False)], size=15, bullet=True, space_after=8)
para(t2, [("And pass-through ", DARK, False), ("λᵢ > 0", ORANGE, True), (" so it reaches farmers.", DARK, False)], size=15, bullet=True, space_after=8)
para(t2, [("Rewards quality along the ", DARK, False), ("continuous gradient", DARK, True), (".", DARK, False)], size=15, bullet=True, space_after=8)
para(t2, [("Exp 1: measurement alone did ", DARK, True), ("not", RED, True), (" activate α.", DARK, True)], size=15, bold=True, space_after=0)
notes(s, "The model formalises why quality can rise with no price response. Two channels. Exp 1 is the enforcement channel "
      "operating in isolation: σ² falls, the floor bites, but α stays at zero because processors still buy flat per litre.")

# =====================================================================
# SLIDE 11 — Model: making the floor bite + CDF
# =====================================================================
s = content_slide("Making the floor bite: the enforcement signature")
img_fit(s, FIG+"fig_cdf_added_water.png", 6.7, 1.7, 6.2, 4.6, align='center', valign='top')
_, tf = add_text(s, 0.6, 1.8, 5.9, 5.0)
para(tf, [("Prediction: ", SLATE, True), ("a credible floor disciplines the lower tail", GREEN, True),
          (", leaving clean suppliers unchanged.", SLATE, True)], size=18, first=True, space_after=14)
para(tf, [("Treatment CDF of added water lies ", DARK, False), ("above", GREEN, True),
          (" control — first-order stochastic dominance.", DARK, False)], size=16, bullet=True, space_after=10)
para(tf, [("Serious adulteration (>5% water): ", DARK, False), ("11% → 5%", RED, True),
          (", roughly halved.", DARK, False)], size=16, bullet=True, space_after=10)
para(tf, [("Median delivery: ", DARK, False), ("unchanged", GREY, True),
          (" — already clean.", DARK, False)], size=16, bullet=True, space_after=10)
para(tf, [("The threat of rejection, once toothless under noise, becomes real. ", SLATE, True),
          ("The floor bites the worst adulterators.", GREEN, True)], size=16, space_after=0)
notes(s, "This CDF is the visual fingerprint of the enforcement channel. The whole treatment effect is in the right tail — "
      "the egregious adulterators — exactly what 'making the floor bite' predicts. Clean suppliers don't move.")

# =====================================================================
# SLIDE 12 — Model -> motivates Exp 2
# =====================================================================
s = content_slide("What the model says is missing")
_, tf = add_text(s, 0.7, 1.8, 12.0, 1.6)
para(tf, [("Enforcement fixed the ", SLATE, False), ("floor", GREEN, True),
          (". To reward quality along the ", SLATE, False), ("gradient", ORANGE, True),
          (", someone must pay a premium — and pass it up.", SLATE, False)], size=21, first=True, space_after=0)
# arrow-style three boxes
b1 = rect(s, 0.7, 3.4, 3.7, 2.4, RGBColor(0xF3,0xF6,0xF2), line=GREY)
t1=b1.text_frame; t1.word_wrap=True; t1.vertical_anchor=MSO_ANCHOR.MIDDLE; t1.margin_left=Pt(12); t1.margin_right=Pt(12)
para(t1, "Can’t move α at the processor", size=16, bold=True, color=SLATE, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(t1, "would need to change how processors price — outside the project’s reach", size=13, color=GREY, align=PP_ALIGN.CENTER, space_after=0)
b2 = rect(s, 4.8, 3.4, 3.7, 2.4, RGBColor(0xFD,0xF0,0xE4), line=ORANGE)
t2=b2.text_frame; t2.word_wrap=True; t2.vertical_anchor=MSO_ANCHOR.MIDDLE; t2.margin_left=Pt(12); t2.margin_right=Pt(12)
para(t2, "So inject λα > 0 one node down", size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(t2, "pay traders directly for verified quality at the trader–MCC interface", size=13, color=DARK, align=PP_ALIGN.CENTER, space_after=0)
b3 = rect(s, 8.9, 3.4, 3.7, 2.4, RGBColor(0xEC,0xF2,0xF8), line=BLUE)
t3=b3.text_frame; t3.word_wrap=True; t3.vertical_anchor=MSO_ANCHOR.MIDDLE; t3.margin_left=Pt(12); t3.margin_right=Pt(12)
para(t3, "Test: is it passed to farmers?", size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(t3, "observe whether traders share the bonus upstream — the pass-through margin", size=13, color=DARK, align=PP_ALIGN.CENTER, space_after=0)
# arrows
for ax in (4.4, 8.5):
    a = add_text(s, ax, 4.25, 0.45, 0.6)[1]
    para(a, "→", size=34, bold=True, color=GREY, align=PP_ALIGN.CENTER, first=True, space_after=0)
notes(s, "Bridge to Exp 2. We can't literally raise α at the processor, so we inject the quality-linked transfer one node "
      "downstream, paying traders for verified quality. This creates a clean pass-through test: do traders share it with farmers?")

# =====================================================================
# SLIDE 13 — Experiment 2 design
# =====================================================================
s = content_slide("Experiment 2: paying traders for quality")
_, tf = add_text(s, 0.7, 1.65, 7.4, 5.3)
para(tf, "Within-MCC RCT on the Exp 1 infrastructure", size=19, bold=True, color=ORANGE, first=True, space_after=12)
para(tf, [("19 MCCs (Kazo) where the analyzer system stayed live; ", DARK, False),
          ("47 treated / 52 control traders", BLUE, True), (", randomised within MCC.", DARK, False)], size=16, bullet=True, space_after=8)
para(tf, [("Both arms tested + told their readings. Treated traders additionally get a ", DARK, False),
          ("cash bonus", ORANGE, True), (" via mobile money: UGX 100/litre per 0.1pp of fat & SNF above standard.", DARK, False)],
     size=16, bullet=True, space_after=8)
para(tf, [("Isolates the ", DARK, False), ("price incentive", ORANGE, True), (" from the information channel.", DARK, False)], size=16, bullet=True, space_after=8)
para(tf, [("Daily mobile-money payment + SMS → ", DARK, False), ("dense daily panel", PURPLE, True),
          (" (5,596 submissions, 31 days).", DARK, False)], size=16, bullet=True, space_after=8)
para(tf, [("Surprise mid-experiment ", DARK, True), ("threshold shift", RED, True), (":", DARK, True)], size=16, bullet=True, space_after=4)
para(tf, [("Stage 1 (fat ≥ 3.3%) ", GREEN, True), ("— easy, rainy season → low q_min.", DARK, False)], size=15, level=1, space_after=2)
para(tf, [("Stage 2 (fat ≥ 3.9%) ", RED, True), ("— many deliveries now fail → high q_min.", DARK, False)], size=15, level=1, space_after=0)
# right info box
rb = rect(s, 8.4, 2.0, 4.3, 3.4, RGBColor(0xFD,0xF0,0xE4), line=None)
t2=rb.text_frame; t2.word_wrap=True; t2.vertical_anchor=MSO_ANCHOR.MIDDLE; t2.margin_left=Pt(14); t2.margin_right=Pt(14)
para(t2, "Why traders?", size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, first=True, space_after=8)
para(t2, "largest unrealised quality gains sit at the trader link", size=14, color=DARK, align=PP_ALIGN.CENTER, space_after=8)
para(t2, "creates a testable pass-through margin to farmers", size=14, color=DARK, align=PP_ALIGN.CENTER, space_after=8)
para(t2, "~100 traders: power without 500 farmers", size=14, color=DARK, align=PP_ALIGN.CENTER, space_after=0)
notes(s, "Design isolates the price channel: both arms get information; only treated get money. Daily payments give a panel. "
      "The unannounced threshold increase moves the regime from intensive-margin (gradient) to extensive-margin (floor) — "
      "the model's key comparative static.")

# =====================================================================
# SLIDE 14 — Exp 2 results: volume response
# =====================================================================
s = content_slide("Experiment 2: traders respond on volume first")
img_fit(s, FIG+"fig_daily_quantity.png", 0.7, 1.7, 7.7, 5.0, align='left', valign='top')
_, tf = add_text(s, 8.7, 2.0, 4.1, 4.8)
para(tf, [("+124 L/day", ORANGE, True), (" delivered by treated traders (~+45%).", DARK, False)], size=18, bullet=True, first=True, space_after=12)
para(tf, [("Response is ", DARK, False), ("immediate", GREEN, True),
          (" and present in ", DARK, False), ("both stages", DARK, True), (".", DARK, False)], size=18, bullet=True, space_after=12)
para(tf, [("When the threshold is easy to clear (Stage 1), the cheapest way to earn more bonus is to ", DARK, False),
          ("source more milk", ORANGE, True), (".", DARK, False)], size=18, bullet=True, space_after=12)
para(tf, [("Volume is a low-cost margin: recruit more suppliers, collect bigger loads.", DARK, False)], size=16, bullet=True, space_after=0)
notes(s, "First result of the daily panel: a large, immediate volume response sustained across both stages. Consistent with "
      "traders working the extensive margin of the bonus when quality is easy — grab volume.")

# =====================================================================
# SLIDE 15 — Exp 2 results: quality only when the floor bites
# =====================================================================
s = content_slide("Quality responds only once the floor bites")
img_fit(s, FIG+"fig_daily_fat.png", 0.7, 1.7, 7.7, 5.0, align='left', valign='top')
_, tf = add_text(s, 8.7, 1.9, 4.1, 5.0)
para(tf, [("Stage 1 (fat ≥ 3.3%): ", GREEN, True), ("no quality response", GREY, True),
          (" (+0.035pp, n.s.).", DARK, False)], size=17, bullet=True, first=True, space_after=12)
para(tf, [("Stage 2 (fat ≥ 3.9%): ", RED, True), ("fat jumps +0.17pp", GREEN, True),
          (" (p<0.001); SNF also rises.", DARK, False)], size=17, bullet=True, space_after=12)
para(tf, [("The ", DARK, False), ("threshold", RED, True),
          (", not the continuous gradient, triggers costly quality sourcing — exactly the model’s extensive-margin prediction.",
          DARK, False)], size=17, bullet=True, space_after=12)
para(tf, [("Quality adjusts ", DARK, False), ("gradually", ORANGE, True),
          (" (~1 week): re-sorting suppliers is costly and lumpy.", DARK, False)], size=17, bullet=True, space_after=0)
notes(s, "Central finding. The continuous gradient — live the whole time — elicits nothing. Only when the discrete "
      "pass/fail floor rises do traders screen for quality. Discrete stakes beat smooth incentives. The ~1-week lag fits "
      "costly supplier re-sorting (and possibly loss aversion).")

# =====================================================================
# SLIDE 16 — Exp 2 results: farmer-level null / negative pass-through
# =====================================================================
s = content_slide("… but farmers are left out: the pass-through fails")
img_fit(s, FIG+"fig_forest_farmer_exp2.png", 0.6, 1.7, 6.6, 4.7, align='left', valign='top')
_, tf = add_text(s, 7.5, 1.9, 5.3, 5.0)
para(tf, [("Farm-gate price ", DARK, False), ("falls", RED, True),
          (": −17 UGX/litre to farmers selling to treated traders (p=0.02).", DARK, False)], size=18, bullet=True, first=True, space_after=12)
para(tf, [("Negative", RED, True), (" pass-through — the opposite of what a price channel should deliver.", DARK, False)], size=18, bullet=True, space_after=12)
para(tf, [("Quality-checking and feeding practices: ", DARK, False), ("essentially null", GREY, True),
          (" (one marginal feeding index, fails multiple-testing correction).", DARK, False)], size=17, bullet=True, space_after=12)
para(tf, [("Trader ", DARK, False), ("endline survey", DARK, True), (": all null too — effects live in the daily panel, not the snapshot.", DARK, False)], size=17, bullet=True, space_after=12)
para(tf, [("Informed intermediaries ", SLATE, True), ("capture the surplus", RED, True),
          (". The premium stops at the trader.", SLATE, True)], size=18, space_after=0)
notes(s, "The pass-through test fails — and worse, prices to farmers fall. Treated traders screen harder and capture the "
      "bonus; farmers see no gain and a price cut. Quality upgrading at the source still has no champion.")

# =====================================================================
# SLIDE 17 — Conclusion
# =====================================================================
s = content_slide("Conclusion: necessary, but not sufficient")
_, tf = add_text(s, 0.7, 1.7, 12.0, 1.4)
para(tf, [("Observability is ", SLATE, True), ("necessary but not sufficient", RED, True),
          (" for a market for quality. Even when prices reward quality, ", SLATE, True),
          ("the structure of intermediation decides who captures the returns.", SLATE, True)],
     size=21, first=True, space_after=0)
_, tf2 = add_text(s, 0.7, 3.15, 12.0, 3.7)
para(tf2, [("Exp 1: ", GREEN, True), ("measurement → enforcement", DARK, True),
           (". Quality rises (less adulteration), but no premium anywhere — processors buy flat per litre.", DARK, False)],
     size=18, bullet=True, first=True, space_after=12)
para(tf2, [("Exp 2: ", ORANGE, True), ("inject a premium → quality rises when the floor binds", DARK, True),
           (", but the surplus is retained by traders and farm-gate prices ", DARK, False), ("fall", RED, True), (".", DARK, False)],
     size=18, bullet=True, space_after=12)
para(tf2, [("Thresholds beat gradients", PURPLE, True),
           (": discrete pass/fail signals move intermediaries facing lumpy adjustment costs; smooth gradients don’t.",
           DARK, False)], size=18, bullet=True, space_after=12)
para(tf2, [("Policy: ", BLUE, True),
           ("measurement + pricing are complementary, not substitutes. Partial fixes leave a low-quality trap intact; "
            "pass-through to farmers needs its own instrument.", DARK, False)], size=18, bullet=True, space_after=0)
notes(s, "Land the thesis: observability necessary, not sufficient. The two experiments together show a low-level "
      "equilibrium trap that single-margin interventions can't escape. Two design lessons: thresholds > gradients for "
      "intermediaries; and pass-through to farmers must be engineered, not assumed. Thank you.")

# =====================================================================
# Remove leftover template slides (original index 0, 2, 3, 4)
# =====================================================================
# Determine xml ids of the original template content slides to drop.
# Original slides were indices 0..4; slide 1 (title) we kept & edited.
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
# The first 5 entries are the original template slides in order; index1 kept.
to_remove = [slides[0], slides[2], slides[3], slides[4]]
for sld in to_remove:
    xml_slides.remove(sld)

# Move the (edited) title slide to the front
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
title_entry = slides[0]  # after removal, original title is now first? verify below
# After removing 0,2,3,4 the remaining order is: [orig1(title), new2, new3, ...]
# title already first -> nothing to do.

prs.save(OUT)
print("Saved", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
print("Slide count:", len(list(prs.slides)))
